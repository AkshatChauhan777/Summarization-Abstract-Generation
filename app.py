import gdown
import streamlit as st
from transformers import pipeline, LEDForConditionalGeneration, LEDTokenizerFast
from langdetect import detect
from googletrans import Translator
import re, os, requests, json
import pandas as pd
from bs4 import BeautifulSoup
import fitz  # PyMuPDF for PDFs
import docx
import pptx
from odf.opendocument import load as load_odt
from odf.text import P
from ebooklib import epub
import xml.etree.ElementTree as ET
import yake
from io import StringIO
import torch

# Constants
WORD_LIMIT = 1600
translator = Translator()

st.title("🧠 AI Based Text Summarization and Abstract Generation")
st.write(f"⚠️ Word limit: {WORD_LIMIT} words.")

# Language detection
def detect_language(text):
    try:
        return detect(text)
    except:
        return "en"

# Safe translation function
def translate_text(text, source_lang, target_lang):
    try:
        if source_lang == target_lang:
            return text
        return translator.translate(text, src=source_lang, dest=target_lang).text
    except Exception as e:
        print(f"Translation failed: {e}")
        return text

# # Load summarization model
# @st.cache_resource
# def load_summarization_model():
#     return pipeline("summarization", model="pszemraj/led-large-book-summary")

# summarizer = load_summarization_model()

# https://drive.google.com/drive/folders/1ft84VHVuVT7LZNzar442JOFqCyf4kLjT?usp=sharing
@st.cache_resource
def download_summarization_model_from_drive():
    model_path = "model1"
    if not os.path.exists(model_path):
        gdown.download_folder(
            "https://drive.google.com/drive/folders/1ft84VHVuVT7LZNzar442JOFqCyf4kLjT?usp=sharing", 
            output=model_path,
            quiet=False,
            use_cookies=False
        )
    tokenizer = LEDTokenizerFast.from_pretrained(model_path)
    model = LEDForConditionalGeneration.from_pretrained(model_path)
    return tokenizer, model

summarization_tokenizer, summarization_model = download_summarization_model_from_drive()


# # Load abstract generation model
# @st.cache_resource
# def load_abstract_model():
#     model_path = "/content/drive/MyDrive/AI Based Summary and Abstract Generator/led-abstract-generator/checkpoint-250"
#     tokenizer = LEDTokenizerFast.from_pretrained(model_path)
#     model = LEDForConditionalGeneration.from_pretrained(model_path)
#     return tokenizer, model

# abstract_tokenizer, abstract_model = load_abstract_model()

import gdown
# https://drive.google.com/drive/folders/1yCdn5sr1Ob5RODeiTi1JCLodoDfYwH-k?usp=sharing
@st.cache_resource
def download_abstract_model_from_drive():
    model_path = "model"
    if not os.path.exists(model_path):
        gdown.download_folder(
            "https://drive.google.com/drive/folders/1yCdn5sr1Ob5RODeiTi1JCLodoDfYwH-k?usp=sharing", quiet=False, use_cookies=False
        )
    tokenizer = LEDTokenizerFast.from_pretrained(model_path)
    model = LEDForConditionalGeneration.from_pretrained(model_path)
    return tokenizer, model

abstract_tokenizer, abstract_model = download_abstract_model_from_drive()


# Clean incomplete sentence
def ensure_complete_sentences(text):
    sentences = re.split(r'(?<=[.!?])\s+', text.strip())
    if len(sentences) > 1:
        return ' '.join(sentences[:-1])
    return text

# Extract text from uploaded file
def extract_text_from_file(uploaded_file, ext):
    text = ""
    if ext == ".txt":
        text = StringIO(uploaded_file.getvalue().decode("utf-8")).read()
    elif ext == ".pdf":
        doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        text = " ".join([page.get_text() for page in doc])
    elif ext == ".docx":
        doc = docx.Document(uploaded_file)
        text = " ".join([para.text for para in doc.paragraphs])
    elif ext in [".xlsx", ".csv"]:
        df = pd.read_excel(uploaded_file) if ext == ".xlsx" else pd.read_csv(uploaded_file)
        text = " ".join(df.astype(str).values.flatten())
    elif ext == ".pptx":
        prs = pptx.Presentation(uploaded_file)
        text = " ".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif ext == ".html":
        soup = BeautifulSoup(uploaded_file.read(), "html.parser")
        text = soup.get_text()
    elif ext == ".odt":
        odt_doc = load_odt(uploaded_file)
        text = " ".join([element.firstChild.data for element in odt_doc.getElementsByType(P) if element.firstChild])
    elif ext == ".epub":
        book = epub.read_epub(uploaded_file)
        text = " ".join([BeautifulSoup(item.content, "html.parser").get_text() for item in book.items if item.get_type() == 9])
    elif ext == ".xml":
        tree = ET.parse(uploaded_file)
        root = tree.getroot()
        text = " ".join(root.itertext())
    elif ext == ".json":
        try:
            data = json.load(uploaded_file)
            text = " ".join(map(str, data.values())) if isinstance(data, dict) else json.dumps(data)
        except Exception as e:
            st.error(f"Error parsing JSON: {e}")
            return ""
    else:
        raise ValueError("Unsupported file format")
    return " ".join(text.split()[:WORD_LIMIT])

# Extract text from URL
def extract_text_from_url(url):
    try:
        response = requests.get(url)
        soup = BeautifulSoup(response.content, "html.parser")
        text = soup.get_text(separator=' ')
        return " ".join(text.split()[:WORD_LIMIT])
    except Exception as e:
        st.error(f"Error extracting text from URL: {e}")
        return ""

# Keyword extraction & highlighting
def extract_keywords_yake(text, max_keywords=15):
    kw_extractor = yake.KeywordExtractor(top=max_keywords, stopwords=None)
    return kw_extractor.extract_keywords(text)

def highlight_keywords(text, keywords):
    keyword_list = [kw for kw, _ in keywords]
    for kw in sorted(keyword_list, key=len, reverse=True):
        escaped_kw = re.escape(kw)
        text = re.sub(fr"\b({escaped_kw})\b", r"**\1**", text, flags=re.IGNORECASE)
    return text

# Summarization core
def summarize_text(text, length_ratio, highlight=True):
    original_lang = detect_language(text)
    translated = False

    if original_lang != "en":
        text = translate_text(text, original_lang, "en")
        translated = True

    word_count = len(text.split())
    min_length = max(10, int(word_count * length_ratio * 0.8))
    max_length = int(word_count * length_ratio * 1.2)

    summary = summarizer(text, min_length=min_length, max_length=max_length, do_sample=False, early_stopping=True)
    summary_text = ensure_complete_sentences(summary[0]['summary_text'])

    if highlight:
        keywords = extract_keywords_yake(summary_text)

        if translated:
            translated_keywords = []
            for kw, score in keywords:
                try:
                    translated_kw = translate_text(kw, "en", original_lang)
                except:
                    translated_kw = kw
                translated_keywords.append((translated_kw, score))
            keywords = translated_keywords

        highlighted_summary = highlight_keywords(summary_text, keywords)

        if translated:
            highlighted_summary = translate_text(highlighted_summary, "en", original_lang)

        return highlighted_summary, keywords, len(summary_text.split())
    else:
        if translated:
            summary_text = translate_text(summary_text, "en", original_lang)
        return summary_text, [], len(summary_text.split())

# Abstract generation logic
def generate_abstract(text, max_input_length=1024, max_output_length=400):
    inputs = abstract_tokenizer(
        text,
        padding="max_length",
        truncation=True,
        max_length=max_input_length,
        return_tensors="pt"
    )
    input_ids = inputs.input_ids
    attention_mask = inputs.attention_mask
    output_ids = abstract_model.generate(
        input_ids=input_ids,
        attention_mask=attention_mask,
        max_length=max_output_length,
        min_length=250,
        num_beams=4,
        repetition_penalty=2.5,
        length_penalty=1.0,
        no_repeat_ngram_size=3,
        early_stopping=True
    )

    return abstract_tokenizer.decode(output_ids[0], skip_special_tokens=True)

# Sidebar: Choose Task
st.sidebar.title("🛠️ Controls Panel")
task = st.sidebar.radio("Choose your task:", ["Generate Summary", "Generate Abstract"])

# ========== SUMMARY GENERATION ==========
if task == "Generate Summary":
    input_method = st.sidebar.radio("Do you want to manually enter text?", ("Yes", "No"))
    text = ""

    if input_method == "Yes":
        text = st.text_area("Enter your text here:", height=300)
        if text:
            st.markdown(f"**🧾 Word Count:** {len(text.split())}")
    else:
        file_option = st.sidebar.radio("Choose Upload Type", ("Upload File", "Enter URL"))

        if file_option == "Upload File":
            uploaded_file = st.sidebar.file_uploader("Upload your document", type=[
                "txt", "pdf", "docx", "xlsx", "csv", "pptx", "html", "odt", "epub", "xml", "json"
            ])
            if uploaded_file:
                file_ext = os.path.splitext(uploaded_file.name)[1].lower()
                try:
                    text = extract_text_from_file(uploaded_file, file_ext)
                    st.success(f"✅ Text extracted successfully from `{uploaded_file.name}`")
                    st.markdown(f"**🧾 Word Count:** {len(text.split())}")
                except Exception as e:
                    st.error(f"❌ Failed to extract text: {e}")
        elif file_option == "Enter URL":
            url = st.sidebar.text_input("Enter URL to extract text:")
            if url:
                text = extract_text_from_url(url)
                if text:
                    st.success("✅ Text extracted from URL")
                    st.markdown(f"**🧾 Word Count:** {len(text.split())}")

    if text:
        st.sidebar.header("Summary Settings")
        summary_type = st.sidebar.selectbox("Choose summary type", ["Concise", "Brief", "Moderate", "Detailed"])
        summary_ratios = {
            "Concise": 0.25,
            "Brief": 0.4,
            "Moderate": 0.55,
            "Detailed": 0.75,
        }

        if st.sidebar.button("Generate Summary"):
            summary, keywords, summary_word_count = summarize_text(text, summary_ratios[summary_type], highlight=True)
            st.markdown(f"### 📝 {summary_type} Summary")
            st.markdown(summary)
            st.markdown(f"**🧾 Summary Word Count:** {summary_word_count}")

            # st.markdown("### 🔑 Keywords")
            # for kw, score in keywords:
            #     st.markdown(f"- {kw} (score: {score:.4f})")

            st.markdown("### 🔑 Keywords")
            for kw, _ in keywords:
                  st.markdown(f"- {kw}")

# ========== ABSTRACT GENERATION ==========
elif task == "Generate Abstract":
    abstract_input = st.text_area("📝 Enter text for abstract generation:", height=300)
    if abstract_input and st.button("Generate Abstract"):
        with st.spinner("Generating abstract..."):
            abstract_text = generate_abstract(abstract_input)
            st.markdown("### 📚 Generated Abstract")
            st.markdown(abstract_text)
            st.markdown(f"**🧾 Abstract Word Count:** {len(abstract_text.split())}")
